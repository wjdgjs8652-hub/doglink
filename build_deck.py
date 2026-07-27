# -*- coding: utf-8 -*-
"""팀 성냥팔아요 — 창업캠프 발표 덱 생성기.
근거: MVP plan.docx / Regional issues report.txt / 권한별기능 보안보강 메모
"""
import math, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

OUT = r"C:\Users\bule1\OneDrive\바탕 화면\team-project\성냥팔아요_발표자료.pptx"

# ── 팔레트: 제주 현무암 + 바다 + 감귤 ────────────────────────────
INK    = RGBColor(0x16, 0x2A, 0x33)   # 현무암 (지배색)
DEEP   = RGBColor(0x0E, 0x1E, 0x25)   # 더 짙은 배경
TEAL   = RGBColor(0x2E, 0x7D, 0x8F)   # 제주 바다 (보조)
TEALL  = RGBColor(0xD9, 0xE8, 0xEC)   # 연한 teal 틴트
ORANGE = RGBColor(0xE0, 0x7A, 0x28)   # 감귤 (액센트)
ORANGL = RGBColor(0xFB, 0xEC, 0xDC)   # 연한 감귤 틴트
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x6B, 0x7A, 0x82)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF6)
LINE   = RGBColor(0xDD, 0xE4, 0xE7)
RED    = RGBColor(0xB4, 0x3D, 0x2E)

HEAD = "맑은 고딕"
BODY = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

WARN = []

def twidth(s, size):
    """대략적 텍스트 폭(pt). 한글은 1em, 라틴/숫자는 0.53em."""
    w = 0.0
    for c in s:
        o = ord(c)
        if o > 0x1100 and not (0x2000 <= o <= 0x206F):
            w += 1.0
        elif c == ' ':
            w += 0.30
        else:
            w += 0.53
    return w * size

def fit(tag, text, box_w_in, box_h_in, size, inset=0.1, lh=1.32):
    """박스 안에 들어가는지 계산으로 확인."""
    avail = (box_w_in - inset * 2) * 72.0
    if avail <= 0:
        WARN.append(f"{tag}: box too narrow"); return
    lines = 0
    for para in text.split("\n"):
        lines += max(1, math.ceil(twidth(para, size) / avail))
    need = lines * size * lh + 6
    have = box_h_in * 72.0
    if need > have:
        WARN.append(f"{tag}: needs {need:.0f}pt, has {have:.0f}pt ({lines} lines @{size}pt)")

def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid()
    bg.fore_color.rgb = DEEP if dark else WHITE
    return s

def tb(s, x, y, w, h, text, size=14, color=INK, bold=False, font=BODY,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, lh=1.32, space_after=0, tag=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    parts = text.split("\n")
    for i, p in enumerate(parts):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = lh
        if space_after:
            para.space_after = Pt(space_after)
        r = para.add_run(); r.text = p
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    if tag:
        fit(tag, text, w, h, size, inset=0.0, lh=lh + (space_after / size if size else 0))
    return box

def card(s, x, y, w, h, fill=LIGHT, line=None, radius=0.06, shadow=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh

def circle(s, x, y, d, fill, label="", color=WHITE, size=15, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = HEAD; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sh

def title(s, text, sub=None, dark=False):
    # sub가 있으면 제목 박스를 1줄 높이로 제한한다 — 2줄로 넘치면 fit()이 잡아낸다
    tb(s, 0.75, 0.62, 11.8, (0.72 if sub else 1.36), text, size=34, bold=True, font=HEAD,
       color=WHITE if dark else INK, tag="title:" + text[:12])
    if sub:
        tb(s, 0.75, 1.44, 11.8, 0.42, sub, size=13.5,
           color=RGBColor(0xB6, 0xC4, 0xCA) if dark else MUTED, tag="sub:" + text[:10])

def kicker(s, text, dark=False):
    tb(s, 0.75, 0.30, 11.8, 0.28, text, size=11.5, bold=True, font=HEAD,
       color=ORANGE if dark else TEAL)

def source(s, text, dark=False):
    tb(s, 0.75, 6.92, 11.8, 0.32, text, size=9.5,
       color=RGBColor(0x8A, 0x99, 0xA0) if dark else MUTED, tag="src")

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def style_chart(ch, colors, val_max=None, fmt='0"%"'):
    ch.has_title = False
    ch.font.name = BODY; ch.font.size = Pt(11); ch.font.color.rgb = INK
    pl = ch.plots[0]
    pl.gap_width = 60
    pl.has_data_labels = True
    dl = pl.data_labels
    dl.font.size = Pt(10.5); dl.font.bold = True; dl.font.name = BODY
    dl.number_format = fmt; dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    for i, ser in enumerate(pl.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        ser.format.line.fill.background()
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = LINE
    va.major_gridlines.format.line.width = Pt(0.75)
    va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = MUTED
    va.tick_labels.font.name = BODY
    va.format.line.fill.background()
    if val_max:
        va.maximum_scale = val_max
        va.minimum_scale = 0
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(11); ca.tick_labels.font.color.rgb = INK
    ca.tick_labels.font.name = BODY
    ca.format.line.color.rgb = LINE

# ════════════════════════════════════════════════════════
# 1. 표지
# ════════════════════════════════════════════════════════
s = slide(dark=True)
circle(s, 0.75, 1.55, 0.62, ORANGE, "", size=1)
tb(s, 0.78, 1.68, 0.6, 0.4, "🐾", size=20, align=PP_ALIGN.CENTER, color=WHITE)
tb(s, 0.75, 2.45, 11.5, 0.5, "2026 제주 지역대학 연합 창업 캠프", size=14, bold=True,
   font=HEAD, color=ORANGE)
tb(s, 0.75, 3.00, 11.5, 1.15, "유기견 발견·제보 연계 서비스", size=52, bold=True,
   font=HEAD, color=WHITE, tag="cover-title")
tb(s, 0.75, 4.28, 10.6, 0.85,
   "QR 한 번으로 1분 안에 끝나는 제보,\n권한에 따라 꼭 필요한 만큼만 열리는 정보",
   size=17, color=RGBColor(0xC3, 0xD2, 0xD8), lh=1.45, tag="cover-sub")
tb(s, 0.75, 5.75, 11.5, 0.32, "팀  성냥팔아요", size=15, bold=True, font=HEAD, color=WHITE)
tb(s, 0.75, 6.14, 11.5, 0.3, "강주원 (기획·총괄)   ·   김 정 (UI/UX·기획)   ·   유정헌 (기획·AI)",
   size=12, color=RGBColor(0x9F, 0xB0, 0xB8))
notes(s, "팀 성냥팔아요입니다. 저희는 제주에서 유기견을 발견한 시민이 1분 안에 제보를 끝내고, "
         "그 정보가 권한에 따라 안전하게 기관까지 전달되는 서비스를 만들었습니다.")

# ════════════════════════════════════════════════════════
# 2. 문제 정의 — 숫자로 보는 제주
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "PROBLEM")
title(s, "제주 유기견 문제는 '늘어나서'가 아니라\n'끝이 나쁘기 때문에' 지역문제입니다")
tb(s, 0.75, 2.12, 11.8, 0.4,
   "구조 건수는 5년째 줄었습니다. 그런데 구조된 뒤의 결과가 전국과 크게 다릅니다.",
   size=14, color=MUTED, tag="p2-lead")

stats = [
    ("39.4%", "감소", "2021→2025 제주 구조견\n4,517 → 2,736마리", TEAL, False),
    ("3.1배", "인구 대비 부담", "전국 평균 대비 제주\n인구 10만 명당 412마리", TEAL, False),
    ("46.7%", "인도적 처리", "전국 17.3%의 약 2.7배\n1,618마리", RED, True),
    ("5.7%", "보호자 반환", "전국 11.2%의 약 절반\n196마리", RED, True),
]
x = 0.75
for val, lab, desc, col, bad in stats:
    card(s, x, 2.66, 2.86, 2.38, fill=(ORANGL if bad else LIGHT))
    tb(s, x + 0.28, 2.94, 2.3, 0.84, val, size=40, bold=True, font=HEAD, color=col, tag="stat" + val)
    tb(s, x + 0.28, 3.80, 2.3, 0.3, lab, size=13, bold=True, font=HEAD, color=INK)
    tb(s, x + 0.28, 4.12, 2.35, 0.78, desc, size=11.5, color=MUTED, lh=1.4, tag="statd" + val)
    x += 3.02

card(s, 0.75, 5.25, 11.83, 1.35, fill=INK)
tb(s, 1.15, 5.52, 11.1, 0.78,
   "\"제주에서 유기견이 계속 늘고 있다\"는 설명은 자료와 맞지 않습니다.\n"
   "저희가 겨냥하는 지점은 발생량이 아니라, 구조된 개가 보호자에게 돌아가지 못하는 비율입니다.",
   size=14.5, color=WHITE, lh=1.5, tag="p2-quote")
source(s, "출처: 농림축산식품부 2025 전국 동일 기준 · 제주동물보호센터 입소 통계 (Regional issues report v2.0 §4·§6·§7)")
notes(s, "먼저 문제를 정확히 짚겠습니다. 제주 구조견은 5년간 39% 줄었습니다. 증가하고 있다고 말하면 틀립니다. "
         "그런데 인도적 처리 비율이 46.7%로 전국의 2.7배, 반환율은 5.7%로 전국의 절반입니다. "
         "즉 문제는 '얼마나 생기느냐'가 아니라 '구조된 뒤 어떻게 끝나느냐'입니다.")

# ════════════════════════════════════════════════════════
# 3. 시장조사
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "MARKET")
title(s, "시장성 판단: 조건부 있음",
      "시민 과금으로는 성립하지 않습니다. 행정 비용 절감과 지역 후원이 실제 지불 주체입니다.")

card(s, 0.75, 2.25, 5.8, 2.05, fill=LIGHT)
circle(s, 1.08, 2.55, 0.42, RED, "✕", size=15)
tb(s, 1.66, 2.62, 4.6, 0.3, "성립하지 않는 시장", size=14, bold=True, font=HEAD, color=INK)
tb(s, 1.08, 3.15, 5.2, 1.0,
   "· 기존 공고 조회·실종 제보 서비스는 대부분 무료\n"
   "· 시민 직접 지불액은 사실상 후원에 가까움\n"
   "· 광고만으로 운영비를 감당하기 어려움",
   size=12.5, color=MUTED, lh=1.55, tag="m1")

card(s, 6.78, 2.25, 5.8, 2.05, fill=TEALL)
circle(s, 7.11, 2.55, 0.42, TEAL, "✓", size=15)
tb(s, 7.69, 2.62, 4.6, 0.3, "성립할 수 있는 시장", size=14, bold=True, font=HEAD, color=INK)
tb(s, 7.11, 3.15, 5.2, 1.0,
   "· 지자체·보호센터의 분산 신고 처리 비용 절감\n"
   "· 기관용 관리 도구로서의 가치\n"
   "· 지역 기업 후원·협업 (미디어, 유통 제품)",
   size=12.5, color=INK, lh=1.55, tag="m2")

tb(s, 0.75, 4.62, 11.8, 0.3, "초기 목표 구조", size=15, bold=True, font=HEAD, color=INK)
items = [
    ("2,000원 / 5,000원", "시민 선택 과금\n광고 제거 / 광고 제거 + 후원"),
    ("월 10~30만 원", "기술 유지비 목표\n서버 · AI 서버 · 스토어 유지"),
    ("청년지원 · 후원", "초기 개발비 조달\n캠프 이후 사업화 재원"),
]
x = 0.75
for big, desc in items:
    card(s, x, 5.05, 3.86, 1.62, fill=WHITE, line=LINE)
    tb(s, x + 0.3, 5.30, 3.3, 0.44, big, size=18, bold=True, font=HEAD, color=ORANGE, tag="mk" + big[:5])
    tb(s, x + 0.3, 5.84, 3.3, 0.7, desc, size=11.5, color=MUTED, lh=1.45, tag="mkd" + big[:5])
    x += 4.03

source(s, "출처: MVP plan.docx §2 시장성 · §3 지불 의사")
notes(s, "시장성은 조건부로 있다고 봅니다. 시민에게 돈을 받는 모델은 성립하지 않습니다. "
         "대신 지자체가 지금 분산된 신고를 처리하며 쓰는 비용을 줄여주는 관리 도구로서 가치가 있습니다. "
         "초기 목표는 매출이 아니라 유효 제보율과 업무 절감 효과를 증명하는 것입니다.")

# ════════════════════════════════════════════════════════
# 4. 제주 유기견 정책 흐름
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "POLICY")
title(s, "제주 정책 흐름: 예방은 이미 작동 중입니다",
      "저희는 예방이 아니라, 예방으로 줄지 않는 '발견 이후' 구간을 맡습니다.")

# 타임라인 (연결선 + 노드)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.06), Inches(2.92), Inches(10.6), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False

tl = [
    ("2019", "읍면지역 실외사육견\n중성화 지원 시작", TEAL),
    ("2022", "중성화 747마리", TEAL),
    ("2023", "동지역까지 사업 확대\n중성화 944마리", TEAL),
    ("2024", "중성화 481마리\n내장형 등록 74.0%", TEAL),
    ("2025", "구조견 2,736마리\n(2021년 대비 -39.4%)", ORANGE),
]
x = 0.75
for yr, desc, col in tl:
    circle(s, x + 0.52, 2.66, 0.54, col, yr[2:], size=13)
    tb(s, x, 3.42, 2.15, 0.3, yr, size=15, bold=True, font=HEAD, color=INK, align=PP_ALIGN.CENTER)
    tb(s, x - 0.06, 3.78, 2.28, 0.85, desc, size=11.5, color=MUTED, lh=1.42,
       align=PP_ALIGN.CENTER, tag="tl" + yr)
    x += 2.43

card(s, 0.75, 4.95, 5.8, 1.72, fill=TEALL)
tb(s, 1.12, 5.22, 5.15, 0.3, "정책이 만든 성과", size=13.5, bold=True, font=HEAD, color=INK)
tb(s, 1.12, 5.62, 5.15, 0.85,
   "내장형 동물등록 74.0%로 전국 평균 51.8%를 22.2%p 앞섭니다.\n"
   "잃어버린 개를 되찾을 기반은 제주가 이미 전국 최고 수준입니다.",
   size=12, color=INK, lh=1.5, tag="pol1")

card(s, 6.78, 4.95, 5.8, 1.72, fill=ORANGL)
tb(s, 7.15, 5.22, 5.15, 0.3, "정책이 닿지 못한 구간", size=13.5, bold=True, font=HEAD, color=INK)
tb(s, 7.15, 5.62, 5.15, 0.85,
   "중성화·등록은 '발생'을 줄입니다. 그러나 이미 길에 나온 개를\n"
   "누가·언제 발견해 어떻게 전달하는지는 여전히 비어 있습니다.",
   size=12, color=INK, lh=1.5, tag="pol2")

source(s, "출처: Regional issues report v2.0 §5 (중성화·동물등록 정책 효과) · §4 (연도별 입소). "
          "감소분 전체가 중성화 때문이라는 단독 인과관계는 미확인.")
notes(s, "제주는 2019년부터 마당개 중성화를 시작해 2023년 동지역까지 넓혔고, 내장형 등록률은 전국 1위 수준입니다. "
         "다만 이 정책들은 '발생'을 줄이는 예방책입니다. 이미 길에 나온 개를 발견해서 기관까지 잇는 구간은 "
         "여전히 비어 있고, 저희는 정확히 그 구간을 맡습니다.")

# ════════════════════════════════════════════════════════
# 5. 남은 격차 (차트)
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "THE GAP")
title(s, "구조 이후 결과: 제주와 전국의 격차",
      "같은 기준(농림축산식품부 2025)으로 비교해도 제주는 반환·입양이 낮고 인도적 처리가 높습니다.")

cd = CategoryChartData()
cd.categories = ["반환", "입양", "자연사", "인도적 처리"]
cd.add_series("제주", (5.7, 16.2, 17.9, 46.7))
cd.add_series("전국", (11.2, 25.6, 25.7, 17.3))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.75), Inches(2.3), Inches(7.5), Inches(4.3), cd)
ch = gf.chart
style_chart(ch, [ORANGE, TEAL], val_max=55, fmt='0.0"%"')
ch.has_legend = True
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11.5)
from pptx.enum.chart import XL_LEGEND_POSITION
ch.legend.position = XL_LEGEND_POSITION.TOP

box = [
    ("2.7배", "인도적 처리 비율\n제주 46.7% vs 전국 17.3%", RED),
    ("64.6%", "자연사 + 인도적 처리 합계\n전국은 43.0%", RED),
    ("절반", "반환 비율\n제주 5.7% vs 전국 11.2%", ORANGE),
]
y = 2.36
for big, desc, col in box:
    card(s, 8.62, y, 3.96, 1.40, fill=LIGHT)
    tb(s, 8.94, y + 0.18, 3.4, 0.58, big, size=25, bold=True, font=HEAD, color=col, tag="g" + big)
    tb(s, 8.94, y + 0.80, 3.4, 0.58, desc, size=11.5, color=MUTED, lh=1.4, tag="gd" + big)
    y += 1.50

source(s, "출처: 농림축산식품부 2025년 전국 동일 기준 조치 현황 (Regional issues report v2.0 §7). "
          "제주센터 관리 기준 보조 수치는 1,778마리·47.9%.")
notes(s, "같은 기준으로 비교한 결과입니다. 제주는 반환이 5.7%로 전국의 절반, 인도적 처리는 46.7%로 2.7배입니다. "
         "자연사까지 합치면 제주는 64.6%가 살아서 나가지 못합니다. "
         "반환율을 끌어올리는 것이 가장 개선 여지가 큰 지점이고, 반환은 '발견 직후 정보'가 결정합니다.")

# ════════════════════════════════════════════════════════
# 6~7. 페르소나
# ════════════════════════════════════════════════════════
def persona(tag, name, meta, scene, quote, pains, needs, accent):
    s = slide()
    kicker(s, tag)
    title(s, name, meta)

    card(s, 0.75, 2.25, 5.8, 4.42, fill=LIGHT)
    circle(s, 1.12, 2.6, 0.72, accent, name[0], size=26)
    tb(s, 2.05, 2.74, 4.2, 0.32, name, size=17, bold=True, font=HEAD, color=INK)
    tb(s, 2.05, 3.10, 4.2, 0.28, meta, size=11.5, color=MUTED)

    tb(s, 1.12, 3.66, 5.1, 0.28, "상황", size=12, bold=True, font=HEAD, color=accent)
    tb(s, 1.12, 4.00, 5.1, 1.20, scene, size=12.5, color=INK, lh=1.5, tag=tag + "-scene")

    tb(s, 1.12, 5.32, 5.1, 0.28, "이 순간의 생각", size=12, bold=True, font=HEAD, color=accent)
    tb(s, 1.12, 5.66, 5.1, 0.92, quote, size=13.5, color=INK, lh=1.5, tag=tag + "-quote")

    tb(s, 6.78, 2.30, 5.8, 0.3, "겪는 불편", size=14, bold=True, font=HEAD, color=INK)
    y = 2.70
    for p in pains:
        circle(s, 6.78, y + 0.01, 0.26, ORANGL, "·", size=13, color=ORANGE)
        tb(s, 7.20, y, 5.35, 0.44, p, size=12.5, color=INK, lh=1.42, tag=tag + "-p")
        y += 0.52

    tb(s, 6.78, y + 0.16, 5.8, 0.3, "그래서 필요한 것", size=14, bold=True, font=HEAD, color=INK)
    y += 0.58
    for n in needs:
        circle(s, 6.78, y + 0.01, 0.26, TEALL, "✓", size=11, color=TEAL)
        tb(s, 7.20, y, 5.35, 0.44, n, size=12.5, color=INK, lh=1.42, tag=tag + "-n")
        y += 0.52
    return s

s = persona(
    "PERSONA A — 주 사용자",
    "이미연", "32세 · 주부 · 제주시 거주 · 산책 중 유기견 발견",
    "아침 산책길에 목줄 없이 혼자 도로 쪽으로 걸어가는 개를 봤습니다.\n"
    "사진은 찍었는데, 그다음에 뭘 해야 할지 모릅니다.",
    "\"신고는 하고 싶은데,\n 어디에 어떻게 해야 하는지 모르겠어요.\"",
    ["어느 기관에 신고해야 하는지 모른다",
     "앱 설치와 회원가입이 부담스럽다",
     "위치를 어떻게 설명해야 할지 어렵다",
     "찾아보는 사이 개가 이동해 버린다"],
    ["설치·가입 없이 즉시 열리는 제보 화면",
     "사진 한 장이면 나머지는 자동으로 정리",
     "내 제보가 어떻게 처리됐는지 확인"],
    TEAL)
notes(s, "주 사용자는 우연히 발견한 시민입니다. 선의는 있지만 신고처를 모르고, 앱 설치 요구를 만나면 이탈합니다. "
         "핵심은 '검색하는 사이 개가 이동한다'는 것 — 시간이 곧 성공률입니다.")

s = persona(
    "PERSONA B — 보조 사용자",
    "강효민", "28세 · 직장인 · 반려견 실종 3일차",
    "퇴근 후 매일 밤 SNS·당근·보호소 공고를 하나씩 다시 확인합니다.\n"
    "목격 제보가 와도 서로 연결되지 않아 어디까지 확인했는지 헷갈립니다.",
    "\"어제 본 그 목격 글이랑\n 오늘 공고가 같은 애인지 모르겠어요.\"",
    ["여러 채널을 반복해서 확인해야 한다",
     "목격 제보끼리 연결되지 않는다",
     "같은 개체인지 판단할 근거가 없다",
     "보호소 처리 상태를 알 수 없다"],
    ["내 실종 신고와 주변 목격의 자동 대조",
     "같은 개체 가능성이 높은 후보 제시",
     "기관 전달·보호 상태 알림"],
    ORANGE)
notes(s, "보조 사용자는 실종 보호자입니다. 정보가 없는 게 아니라 흩어져 있는 게 문제입니다. "
         "여러 채널을 반복 확인하면서도 같은 개체인지 판단할 방법이 없습니다.")

# ════════════════════════════════════════════════════════
# 8. 사용자 조사 — 제보 불이행 원인
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "USER RESEARCH")
title(s, "제보가 행동으로 이어지지 않는 5가지 이유",
      "선의가 없어서가 아니라, 선의와 행동 사이의 마찰 때문입니다.")

causes = [
    ("01", "신고처를 모른다", "어디에 알려야 하는지 몰라\n검색만 반복한다", TEAL),
    ("02", "설치·가입 장벽", "앱 설치와 회원가입 요구를\n만나는 순간 이탈한다", TEAL),
    ("03", "입력이 어렵다", "위치를 어떻게 특정할지\n판단이 서지 않는다", TEAL),
    ("04", "정보가 비정형", "채널마다 항목이 달라\n누락과 중복이 남는다", ORANGE),
    ("05", "결과를 모른다", "구조·반환 결과를 몰라\n다시 할 이유가 없다", ORANGE),
]
x = 0.75
for num, head, desc, col in causes:
    card(s, x, 2.28, 2.28, 2.46, fill=WHITE, line=LINE)
    circle(s, x + 0.28, 2.56, 0.46, col, num, size=13)
    tb(s, x + 0.28, 3.22, 1.9, 0.58, head, size=13.5, bold=True, font=HEAD, color=INK,
       lh=1.25, tag="c" + num)
    tb(s, x + 0.28, 3.88, 1.92, 0.78, desc, size=10.5, color=MUTED, lh=1.4, tag="cd" + num)
    x += 2.4

card(s, 0.75, 5.0, 11.83, 1.62, fill=ORANGL)
circle(s, 1.15, 5.28, 0.44, ORANGE, "!", size=17)
tb(s, 1.78, 5.28, 10.4, 0.3, "이 5가지는 아직 '검증된 조사 결과'가 아니라 가설입니다",
   size=14, bold=True, font=HEAD, color=INK)
tb(s, 1.78, 5.68, 10.4, 0.82,
   "근거자료 12-2절은 시민의 미인지 비율·중도 이탈률·기관 확인까지 걸리는 시간을 '아직 확인되지 않은 항목'으로 분류합니다.\n"
   "캠프 이후 시민 20명 인터뷰, 기관 5곳 인터뷰, 실제 신고 사례 30건 추적으로 검증한 뒤 확정하겠습니다.",
   size=12, color=INK, lh=1.5, tag="ur-warn")

source(s, "가설 근거: MVP plan.docx §4 (문제 상황·현재 방식의 한계) · Regional issues report v2.0 §12-2, §12-4")
notes(s, "여기가 가장 솔직하게 말씀드려야 할 부분입니다. 이 5가지 원인은 저희가 문서와 기존 자료에서 도출한 가설이지, "
         "설문으로 검증한 숫자가 아닙니다. 근거 보고서 자체가 이 항목들을 '아직 확인되지 않았다'고 명시합니다. "
         "그래서 저희는 캠프 이후 검증 계획을 함께 제시합니다.")

# ════════════════════════════════════════════════════════
# 9. 핵심 가치
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "CORE VALUE")
title(s, "저희가 제공하는 네 가지 가치",
      "앞의 다섯 가지 마찰을 하나씩 제거하는 방식으로 설계했습니다.")

vals = [
    ("1분", "설치도 가입도 없는 제보",
     "웹 포스터 QR을 찍으면 바로 제보 화면이 열립니다.\n사진·위치·목격 내용만 넣으면 끝납니다.", TEAL),
    ("표준화", "AI가 항목을 정리합니다",
     "자유롭게 쓴 문장에서 털색·체격·무늬·부상·이동 방향을\n뽑아 표준 카드로 만들고, 빠진 항목을 짚어 줍니다.", TEAL),
    ("가시화", "처리 상태가 보입니다",
     "제보됨 → 확인 중 → 기관 전달 → 보호/반환.\n각 단계에 머문 시간까지 계측해 병목을 드러냅니다.", ORANGE),
    ("최소 공개", "권한만큼만 열립니다",
     "시민 응답에는 정확한 좌표를 아예 담지 않습니다.\n기관은 사유를 남겨야 좌표를 열 수 있습니다.", ORANGE),
]
positions = [(0.75, 2.28), (6.78, 2.28), (0.75, 4.55), (6.78, 4.55)]
for (vx, vy), (big, head, desc, col) in zip(positions, vals):
    card(s, vx, vy, 5.8, 2.05, fill=LIGHT)
    circle(s, vx + 0.34, vy + 0.32, 0.92, col, big, size=15)
    tb(s, vx + 1.48, vy + 0.36, 4.05, 0.38, head, size=15, bold=True, font=HEAD, color=INK,
       tag="v" + big)
    tb(s, vx + 1.48, vy + 0.84, 4.05, 0.95, desc, size=12, color=MUTED, lh=1.5, tag="vd" + big)

source(s, "구현 상태: 네 가지 모두 작동하는 프로토타입으로 시연 가능 (mvp-prototype.html)")
notes(s, "핵심 가치는 네 가지입니다. 1분 안에 끝나는 무설치 제보, AI 표준화, 처리 상태 가시화, 그리고 권한별 최소 공개입니다. "
         "중요한 건 이 네 가지가 슬라이드 위의 약속이 아니라 지금 열어서 보여드릴 수 있는 프로토타입이라는 점입니다.")

# ════════════════════════════════════════════════════════
# 10. 경쟁사 조사
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "COMPETITION")
title(s, "경쟁 환경: 가장 강한 경쟁자는 이미 존재하는 국가 시스템입니다")

rows = [
    ("국가동물보호정보시스템", "정부", "동물 발견 신고, 담당자 문자 자동 발송,\n분실 게시물, 보호 동물 조회",
     "설치 없이 쓸 수 있으나 인지도·완주율 미확인", RED),
    ("포인핸드", "민간 앱", "보호소 공고 조회, 실종·보호 제보",
     "앱 설치 필요, 처리 상태 연동 없음", MUTED),
    ("하트퍼즈", "민간 앱", "위치 기반 실종 SOS",
     "실종 보호자 중심, 발견 시민 동선은 약함", MUTED),
    ("PawBoost", "해외", "웹 제보와 지역 확산",
     "국내 기관 처리 절차와 연결되지 않음", MUTED),
]
hy = 2.16
tb(s, 0.95, hy, 3.0, 0.28, "서비스", size=11, bold=True, font=HEAD, color=MUTED)
tb(s, 3.95, hy, 1.1, 0.28, "구분", size=11, bold=True, font=HEAD, color=MUTED)
tb(s, 5.15, hy, 3.6, 0.28, "제공 기능", size=11, bold=True, font=HEAD, color=MUTED)
tb(s, 8.95, hy, 3.6, 0.28, "한계 / 우리와의 접점", size=11, bold=True, font=HEAD, color=MUTED)

y = 2.52
for name, kind, feat, gap, col in rows:
    card(s, 0.75, y, 11.83, 0.88, fill=(ORANGL if col == RED else LIGHT))
    tb(s, 0.95, y + 0.15, 2.95, 0.6, name, size=12.5, bold=True, font=HEAD, color=INK,
       lh=1.25, tag="cm" + name[:6])
    tb(s, 3.95, y + 0.29, 1.1, 0.3, kind, size=11, color=MUTED)
    tb(s, 5.15, y + 0.15, 3.65, 0.62, feat, size=11, color=INK, lh=1.35, tag="cf" + name[:6])
    tb(s, 8.95, y + 0.15, 3.5, 0.62, gap, size=11, color=(RED if col == RED else MUTED),
       lh=1.35, tag="cg" + name[:6])
    y += 0.96

card(s, 0.75, 6.40, 11.83, 0.42, fill=TEALL)
tb(s, 1.05, 6.51, 11.3, 0.30,
   "차별점 — 별도 창구를 또 만들지 않고, QR 진입 · AI 표준화 · 상태 연결 · 권한별 최소 공개를 하나의 흐름으로 잇습니다.",
   size=11.5, bold=True, font=HEAD, color=INK, tag="cmp-diff")
source(s, "출처: MVP plan.docx §2 유사 기업 · Regional issues report v2.0 §12-1 (기존 시스템 기능 확인)")
notes(s, "정직하게 말씀드리면, 저희의 가장 큰 경쟁자는 다른 스타트업이 아니라 이미 있는 국가동물보호정보시스템입니다. "
         "'온라인 신고 방법이 없다'는 주장은 사실이 아닙니다. 다만 그 시스템의 인지도와 완주율은 확인된 자료가 없고, "
         "저희는 별도 창구를 또 만드는 대신 진입·표준화·상태 연결을 잇는 방향으로 차별화합니다.")

# ════════════════════════════════════════════════════════
# 11. SWOT
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "SWOT")
title(s, "SWOT 분석")

sw = [
    ("S", "강점", TEAL, WHITE, [
        "작동하는 프로토타입 보유 (기획서가 아닌 시연물)",
        "QR 무설치 진입 — 설치·가입 장벽 제거",
        "권한별 데이터 최소 공개를 설계 단계에서 반영",
        "제주 한 지역에서 작게 검증 가능한 범위"]),
    ("W", "약점", ORANGE, WHITE, [
        "실사용자 검증 0건 — 원인 진단이 아직 가설",
        "기관 협조·데이터 연동 미확보",
        "AI는 현재 규칙 기반 모의 응답 (비전 모델 미연결)",
        "시민 과금 모델의 수익성이 낮음"]),
    ("O", "기회", TEAL, WHITE, [
        "인도적 처리 46.7% — 개선 여지가 큰 지표",
        "내장형 등록 74.0%로 반환 인프라는 이미 우수",
        "들개 피해 158건 등으로 행정·주민 관심 높음",
        "지자체 분산 신고 처리 비용 절감 수요"]),
    ("T", "위협", RED, WHITE, [
        "국가동물보호정보시스템과 기능 중복 위험",
        "병목이 신고가 아니라 포획 인력·보호 공간일 가능성",
        "위치정보가 불법 포획·학대에 악용될 위험",
        "개인정보·위치정보 규제 대응 부담"]),
]
positions = [(0.75, 2.12), (6.78, 2.12), (0.75, 4.48), (6.78, 4.48)]
for (vx, vy), (letter, label, col, tcol, items) in zip(positions, sw):
    card(s, vx, vy, 5.8, 2.18, fill=LIGHT)
    circle(s, vx + 0.3, vy + 0.26, 0.5, col, letter, size=17)
    tb(s, vx + 0.94, vy + 0.36, 3.0, 0.3, label, size=14.5, bold=True, font=HEAD, color=INK)
    ly = vy + 0.86
    for it in items:
        tb(s, vx + 0.32, ly, 5.2, 0.3, "· " + it, size=11, color=INK, lh=1.3, tag="sw" + letter)
        ly += 0.32

source(s, "위협 2번은 근거자료 12-4의 '서비스 보류 기준'입니다 — 병목이 신고가 아님이 확인되면 별도 앱 개발을 우선하지 않습니다.")
notes(s, "SWOT입니다. 약점과 위협을 특히 솔직하게 적었습니다. "
         "가장 큰 위협은 근거 보고서가 직접 제시한 보류 기준입니다. 문제의 병목이 신고가 아니라 포획 인력과 보호 공간이라면, "
         "저희는 별도 앱 개발을 우선하지 않아야 합니다. 그래서 다음 슬라이드의 계측 기능이 중요합니다.")

# ════════════════════════════════════════════════════════
# 12. 솔루션 / MVP
# ════════════════════════════════════════════════════════
s = slide()
kicker(s, "SOLUTION")
title(s, "MVP — 핵심 기능 3가지, 지금 작동합니다",
      "QR 포스터를 찍는 순간부터 기관이 상태를 바꾸는 순간까지 하나의 흐름으로 시연합니다.")

feats = [
    ("Core", "QR 기반 AI 보조 무설치 제보",
     "QR 진입 → 사진·위치·목격 내용 입력 →\nAI 표준 제보 카드 생성 → 사용자 확정 →\n접수번호 발급까지 1분 이내", TEAL),
    ("Supporting 1", "제보 지도 · 처리 상태 · 시간 계측",
     "위험도별 공개 범위, 단계별 소요시간 계측,\n지연 사건 표시. 시민에게는 범위만,\n기관에는 사유를 남긴 뒤 정확 좌표", ORANGE),
    ("Supporting 2", "중복 · 실종 후보 매칭",
     "털색·무늬·체격·거리·시간차를 비교해\n후보 최대 3건과 근거 제시.\nAI는 확정하지 않고 사람이 결정", TEAL),
]
x = 0.75
for kind, head, desc, col in feats:
    card(s, x, 2.32, 3.86, 2.5, fill=WHITE, line=LINE)
    tb(s, x + 0.32, 2.6, 3.2, 0.26, kind, size=10.5, bold=True, font=HEAD, color=col)
    tb(s, x + 0.32, 2.94, 3.25, 0.62, head, size=14.5, bold=True, font=HEAD, color=INK,
       lh=1.25, tag="f" + kind)
    tb(s, x + 0.32, 3.66, 3.3, 1.0, desc, size=11.5, color=MUTED, lh=1.45, tag="fd" + kind)
    x += 4.03

card(s, 0.75, 5.08, 11.83, 1.55, fill=INK)
tb(s, 1.15, 5.3, 4.6, 0.3, "설계 원칙", size=12.5, bold=True, font=HEAD, color=ORANGE)
tb(s, 1.15, 5.68, 11.0, 0.8,
   "\"일반 사용자 화면에서 위치를 가리는 것만으로는 부족합니다.\n"
   " 일반 사용자에게 전송되는 데이터 자체에서 정확한 좌표를 제거해야 합니다.\"",
   size=14, color=WHITE, lh=1.5, tag="sol-q")

source(s, "권한별 응답 필드 수: 일반 시민 13개 / 기관 담당자 19개 — 시민 응답에는 좌표·연락처·내부 메모가 존재하지 않음")
notes(s, "MVP는 문서가 정한 핵심 기능 3가지 그대로입니다. 셋 다 지금 브라우저에서 시연 가능합니다. "
         "설계 원칙은 팀 보안 메모에서 가져왔습니다. 화면에서 가리는 게 아니라 응답 데이터에서 지웁니다. "
         "실제로 시민 권한은 13개 필드, 기관은 19개 필드를 받고, 시민 응답에는 좌표가 아예 없습니다.")

# ════════════════════════════════════════════════════════
# 13. 검증 계획 / 마무리
# ════════════════════════════════════════════════════════
s = slide(dark=True)
kicker(s, "NEXT", dark=True)
title(s, "다음 단계 — 주장 대신 데이터를 만듭니다", dark=True)

steps = [
    ("01", "사용자 검증", "시민·발견 경험자 20명 인터뷰\n국가시스템 사용성 테스트 30명"),
    ("02", "기관 검증", "제주동물보호센터·제주시·서귀포시\n자치경찰 등 5개 기관 인터뷰"),
    ("03", "사례 추적", "실제 신고 사례 30건 처리 흐름\n야간·주말·읍면 대응 절차 확인"),
    ("04", "지표 확보", "신고 완료율, 최초 확인 시간,\n반환까지 걸린 시간, 중복 비율"),
]
x = 0.75
for num, head, desc in steps:
    card(s, x, 2.35, 2.86, 2.25, fill=RGBColor(0x1B, 0x2E, 0x37))
    circle(s, x + 0.3, 2.62, 0.46, ORANGE, num, size=13)
    tb(s, x + 0.3, 3.26, 2.3, 0.32, head, size=13.5, bold=True, font=HEAD, color=WHITE)
    tb(s, x + 0.3, 3.68, 2.4, 0.78, desc, size=11, color=RGBColor(0xA8, 0xBA, 0xC2), lh=1.42,
       tag="n" + num)
    x += 3.02

tb(s, 0.75, 4.95, 11.83, 0.9,
   "추진 기준 — 신고 완료 시간 또는 담당 기관 연결 시간을 의미 있게 줄일 수 있을 때 추진합니다.\n"
   "보류 기준 — 병목이 포획 인력·보호 공간·입양 수요 부족으로 확인되면 별도 앱 개발을 우선하지 않습니다.",
   size=13, color=RGBColor(0xC3, 0xD2, 0xD8), lh=1.55, tag="crit")

tb(s, 0.75, 5.92, 11.83, 0.88,
   "시민에게는 구조에 참여할 수 있는 정보만,\n인증된 기관에는 구조를 실행할 수 있는 정보만 제공합니다.",
   size=19, bold=True, font=HEAD, color=WHITE, lh=1.4, tag="closing")

source(s, "검증 항목 출처: Regional issues report v2.0 §12-4 (개발 전 필수 검증)", dark=True)
notes(s, "마무리입니다. 저희는 캠프에서 '이 서비스가 반드시 필요하다'고 주장하는 대신, "
         "그 주장을 검증할 계획과 그 데이터를 수집하는 도구를 함께 내놓습니다. "
         "추진 기준과 보류 기준을 모두 명시했습니다. 데이터가 아니라고 말하면 저희는 방향을 바꿉니다.")

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
print("WARNINGS:", len(WARN))
for w in WARN:
    print("  -", w)



