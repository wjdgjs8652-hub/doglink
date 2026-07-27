# -*- coding: utf-8 -*-
"""렌더링 없이 하는 기하 QA: 텍스트 겹침, 슬라이드 밖 이탈, 여백 부족."""
import io, sys
from pptx import Presentation
from pptx.util import Emu

P = r"C:\Users\bule1\OneDrive\바탕 화면\team-project\성냥팔아요_발표자료.pptx"
prs = Presentation(P)
SW = prs.slide_width / 914400.0
SH = prs.slide_height / 914400.0
MARGIN = 0.5

out = io.open(r"C:\Users\bule1\AppData\Local\Temp\claude\C--Users-bule1-OneDrive-------team-project\485d88a6-63fd-4f78-9230-4a767b2c67de\scratchpad\qa_report.txt",
              "w", encoding="utf-8")

def box(sh):
    return (sh.left/914400.0, sh.top/914400.0,
            (sh.left+sh.width)/914400.0, (sh.top+sh.height)/914400.0)

def overlap(a, b):
    ox = min(a[2], b[2]) - max(a[0], b[0])
    oy = min(a[3], b[3]) - max(a[1], b[1])
    return ox, oy

problems = 0
for i, sl in enumerate(prs.slides, 1):
    texts = []      # (name, box, text)
    for sh in sl.shapes:
        if sh.shape_type is not None and sh.has_text_frame and sh.text_frame.text.strip():
            kind = "TXT" if sh.shape_type == 17 else "SHP"   # 17 = TEXT_BOX
            texts.append((kind, box(sh), sh.text_frame.text.strip().replace("\n", " / ")))
        b = box(sh)
        if b[0] < -0.01 or b[1] < -0.01 or b[2] > SW + 0.01 or b[3] > SH + 0.01:
            out.write(f"[S{i}] OUT-OF-BOUNDS {b} :: {getattr(sh,'name','?')}\n"); problems += 1

    tonly = [t for t in texts if t[0] == "TXT"]
    for a in range(len(tonly)):
        for b_ in range(a+1, len(tonly)):
            ox, oy = overlap(tonly[a][1], tonly[b_][1])
            if ox > 0.02 and oy > 0.02:
                out.write(f"[S{i}] TEXT-OVERLAP {ox:.2f}x{oy:.2f}in\n"
                          f"        A: {tonly[a][2][:52]}\n        B: {tonly[b_][2][:52]}\n")
                problems += 1
    # 텍스트가 슬라이드 여백 안으로 들어왔는지
    for k, bx, tx in tonly:
        if bx[0] < MARGIN - 0.01 or bx[2] > SW - MARGIN + 0.01:
            out.write(f"[S{i}] MARGIN x=({bx[0]:.2f},{bx[2]:.2f}) :: {tx[:44]}\n"); problems += 1

out.write(f"\nTOTAL PROBLEMS: {problems}\n")

# 내용 덤프
out.write("\n" + "="*60 + "\nCONTENT DUMP\n" + "="*60 + "\n")
for i, sl in enumerate(prs.slides, 1):
    out.write(f"\n--- Slide {i} ---\n")
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.write("  " + sh.text_frame.text.strip().replace("\n", "\n  ") + "\n")
        if sh.has_chart:
            ch = sh.chart
            cats = [str(c) for c in ch.plots[0].categories]
            out.write(f"  [CHART] cats={cats}\n")
            for ser in ch.plots[0].series:
                out.write(f"          {ser.name}: {list(ser.values)}\n")
    if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip():
        out.write("  [NOTES] " + sl.notes_slide.notes_text_frame.text.strip()[:90] + "...\n")
out.close()
print("problems:", problems)
